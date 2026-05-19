from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(14, 17, 23)       # #0E1117
CYAN = RGBColor(0, 229, 255)         # #00E5FF
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(200, 200, 200)
AMBER = RGBColor(255, 214, 0)

# Helper function to add a slide with dark background
def add_dark_slide():
    slide_layout = prs.slide_layouts[6] # blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add dark background rectangle
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_BG
    rect.line.fill.background()
    return slide

# Helper function to add standard titles
def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CYAN
    return title_box

# Slide 1: Title Slide
slide1 = add_dark_slide()
title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(2.8))
tf = title_box.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]
p1.text = "AI-POWERED AUTONOMOUS\nCROWD MONITORING & IoT GATEWAY"
p1.font.name = 'Arial'
p1.font.size = Pt(42)
p1.font.bold = True
p1.font.color.rgb = CYAN
p1.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = "Real-Time Edge Intelligence, Computer Vision, & Proactive Safety Actuation"
p2.font.name = 'Arial'
p2.font.size = Pt(18)
p2.font.color.rgb = GRAY
p2.space_before = Pt(15)

author_box = slide1.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(11.83), Inches(1.5))
tf_a = author_box.text_frame
tf_a.word_wrap = True
p_a = tf_a.paragraphs[0]
p_a.text = "Presenter: Manoj Gowda\nDomain: Intelligent Edge Systems & IoT Automation\nEvent: Tech Expo 2026"
p_a.font.name = 'Arial'
p_a.font.size = Pt(14)
p_a.font.color.rgb = WHITE

# Slide 2: Problem Statement
slide2 = add_dark_slide()
add_title(slide2, "THE CHALLENGE: PUBLIC HAZARDS & PASSIVE SYSTEMS")

col1 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.8))
tf1 = col1.text_frame
tf1.word_wrap = True

points_left = [
    ("The Crowd Hazard", "Rapid urban clustering, massive transit hubs, and commercial gatherings scale physical stampede risks exponentially."),
    ("Passive CCTV Limitations", "Traditional CCTV cameras act purely as standard 'witnesses'—they record incidents instead of actively preventing them.")
]
for title, desc in points_left:
    p_t = tf1.add_paragraph() if tf1.paragraphs[0].text else tf1.paragraphs[0]
    p_t.text = "🔴 " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = AMBER
    p_t.space_before = Pt(12)
    
    p_d = tf1.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

col2 = slide2.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
tf2 = col2.text_frame
tf2.word_wrap = True

points_right = [
    ("The Human Guard Bottleneck", "Guards experience severe cognitive fatigue. Watching dozens of feeds leads to a 90% drop in active observation after 20 minutes."),
    ("Disconnected Infrastructure", "Traditional camera analytics do not integrate with building hardware, leaving exits, locks, and sirens dumb and manual.")
]
for title, desc in points_right:
    p_t = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p_t.text = "🔴 " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = AMBER
    p_t.space_before = Pt(12)
    
    p_d = tf2.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

# Slide 3: Proposed Solution
slide3 = add_dark_slide()
add_title(slide3, "THE SOLUTION: ACTIVE EDGE INTELLIGENCE GATEWAY")

col1 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.8))
tf1 = col1.text_frame
tf1.word_wrap = True

sol_left = [
    ("Active AI Architecture", "Shifts security cameras from mere recorders to proactive system decision-makers controlling physical triggers in real-time."),
    ("Edge-First Operations", "Runs full deep-learning algorithms locally on local servers, eliminating latency and securing continuous offline function.")
]
for title, desc in sol_left:
    p_t = tf1.add_paragraph() if tf1.paragraphs[0].text else tf1.paragraphs[0]
    p_t.text = "⚡ " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN
    p_t.space_before = Pt(12)
    
    p_d = tf1.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

col2 = slide3.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
tf2 = col2.text_frame
tf2.word_wrap = True

sol_right = [
    ("Unified Hardware Integration", "Connects deep AI insights with standard physical devices (alarms, automatic gates) using standard lightweight messaging."),
    ("Automated Safety Actuation", "Instantaneous alarm triggers and smart magnetic lock deployments during peak crowd density or panic detections.")
]
for title, desc in sol_right:
    p_t = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p_t.text = "⚡ " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN
    p_t.space_before = Pt(12)
    
    p_d = tf2.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

# Slide 4: System Architecture / Pipeline
slide4 = add_dark_slide()
add_title(slide4, "SYSTEM PIPELINE & MODULAR DATA FLOW")

flow_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.2))
tf = flow_box.text_frame
tf.word_wrap = True

steps = [
    ("1. High-Speed Visual Capture", "Captures raw high-definition video frames from standard surveillance cameras or local webcam devices."),
    ("2. Real-Time Tracking Core", "YOLOv8 Edge AI detects objects, and the ByteTrack algorithm maps unique IDs persistently to tracks."),
    ("3. Local Relational Storage", "Database logs crowd thresholds and manages user security credentials with robust Bcrypt hashing."),
    ("4. RESTful API & MQTT Broker", "Flask exposes live status JSON endpoints, while Paho-MQTT publishes instant telemetry signals."),
    ("5. Autonomous Actuation", "External IoT devices (buzzers, smart door relays) subscribe to MQTT, engaging instantly during triggers.")
]

for title, desc in steps:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_before = Pt(14)
    
    run1 = p.add_run()
    run1.text = title + "  ➔  "
    run1.font.name = 'Arial'
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = CYAN
    
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = 'Arial'
    run2.font.size = Pt(15)
    run2.font.bold = False
    run2.font.color.rgb = WHITE

# Slide 5: Core Tech Stack
slide5 = add_dark_slide()
add_title(slide5, "CORE TECH STACK: STABLE, SOTA UTILITIES")

quads = [
    ("ARTIFICIAL INTELLIGENCE", ["YOLOv8 Object Detection (Ultra-fast, accurate)", "ByteTrack Multi-Object ID Persistence", "Homography Transformation (2D to 3D mapping)"], Inches(0.75), Inches(1.8)),
    ("INTEGRATED NETWORKING", ["Paho-MQTT Protocol (Ultra-low latency telemetry)", "Flask Server (Headless RESTful API Gateway)", "JSON-Based Data Payload standard"], Inches(6.9), Inches(1.8)),
    ("CORE UTILITIES & STORAGE", ["Bcrypt Library (One-way hashed password credentials)", "SQLite relational database (Offline edge file)", "NumPy & Math (Matrix operations)"], Inches(0.75), Inches(4.5)),
    ("ANALYTICS & VISUALS", ["Streamlit Dashboard (Interactive control hub)", "Pandas DataFrames (Live analytical engine)", "Altair Vector Plots (Interactive trend curves)"], Inches(6.9), Inches(4.5))
]

for label, items, x, y in quads:
    box = slide5.shapes.add_textbox(x, y, Inches(5.6), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔳 " + label
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    for item in items:
        pi = tf.add_paragraph()
        pi.text = "• " + item
        pi.font.name = 'Arial'
        pi.font.size = Pt(13)
        pi.font.color.rgb = WHITE
        pi.space_before = Pt(3)

# Slide 6: Dynamic Dashboard Features
slide6 = add_dark_slide()
add_title(slide6, "THE DYNAMIC INTERACTIVE SECURITY HUB")

col1 = slide6.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.8))
tf1 = col1.text_frame
tf1.word_wrap = True

features_left = [
    ("3D Bird's-Eye Radar View", "Translates raw 2D camera coordinates into a flat 3D coordinate model, giving real-time grid positioning of crowd density."),
    ("Zone Boundary Breaches", "Define virtual restricted gates. The system automatically alerts and logs timestamps of unauthorized border breaches.")
]
for title, desc in features_left:
    p_t = tf1.add_paragraph() if tf1.paragraphs[0].text else tf1.paragraphs[0]
    p_t.text = "🎛️ " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = AMBER
    p_t.space_before = Pt(12)
    
    p_d = tf1.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

col2 = slide6.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
tf2 = col2.text_frame
tf2.word_wrap = True

features_right = [
    ("Panic & Anomaly Tracking", "Computes relative velocities of tracked elements. Sudden speed spikes flag dynamic running/panic instances instantly."),
    ("AI Gender Estimation", "Supports high-speed gender analysis simulation, logging visitor demographics for commercial operations.")
]
for title, desc in features_right:
    p_t = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p_t.text = "🎛️ " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = AMBER
    p_t.space_before = Pt(12)
    
    p_d = tf2.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

# Slide 7: Applications & Impact
slide7 = add_dark_slide()
add_title(slide7, "PRACTICAL APPLICATIONS & PUBLIC DEPLOYMENT")

col1 = slide7.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.8))
tf1 = col1.text_frame
tf1.word_wrap = True

apps_left = [
    ("Public Transit Infrastructure", "Platform overcrowding detection at railway and metro terminals, linking alerts directly to gate access and traffic controls."),
    ("Stadiums & Mass Events", "Tracks bottlenecking points at exits and food courts, automating queue routing to prevent crush conditions.")
]
for title, desc in apps_left:
    p_t = tf1.add_paragraph() if tf1.paragraphs[0].text else tf1.paragraphs[0]
    p_t.text = "🏢 " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN
    p_t.space_before = Pt(12)
    
    p_d = tf1.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

col2 = slide7.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
tf2 = col2.text_frame
tf2.word_wrap = True

apps_right = [
    ("Active Emergency Routing", "Computes BFS escape paths in real-time, avoiding dense clusters and lighting smart exit paths dynamically."),
    ("Smart Buildings & Retail Analytics", "Logs detailed customer traffic flow, demographic distributions, and loitering warnings in restricted halls.")
]
for title, desc in apps_right:
    p_t = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p_t.text = "🏢 " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN
    p_t.space_before = Pt(12)
    
    p_d = tf2.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

# Slide 8: Future Scope
slide8 = add_dark_slide()
add_title(slide8, "FUTURE ENHANCEMENTS: BEHAVIORAL DE-ESCALATION")

col1 = slide8.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.8))
tf1 = col1.text_frame
tf1.word_wrap = True

future_left = [
    ("Behavioral De-escalation IoT", "Leverages ambient smart lighting (soothing blue/green hues) and automated calming speech prompts to de-escalate crowd panic at triggers."),
    ("Predictive Trend Forecasting", "Integrates time-series deep neural networks (LSTM) to evaluate logged database history, forecasting overcrowding hours prior.")
]
for title, desc in future_left:
    p_t = tf1.add_paragraph() if tf1.paragraphs[0].text else tf1.paragraphs[0]
    p_t.text = "🚀 " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN
    p_t.space_before = Pt(12)
    
    p_d = tf1.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

col2 = slide8.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
tf2 = col2.text_frame
tf2.word_wrap = True

future_right = [
    ("Dynamic Drone Integration", "Deploys remote multi-camera networks with 5G Network Slicing to guarantee instant, high-speed telemetry in remote environments."),
    ("Blockchain Security Ledgers", "Logs alerts to a decentralized blockchain ledger, offering unalterable, legally compliant data auditing.")
]
for title, desc in future_right:
    p_t = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p_t.text = "🚀 " + title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN
    p_t.space_before = Pt(12)
    
    p_d = tf2.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = WHITE
    p_d.space_before = Pt(4)

prs.save("Crowd_Monitoring_Presentation.pptx")
print("Presentation generated successfully!")
